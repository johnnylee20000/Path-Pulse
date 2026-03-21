"""
Build Path-Pulse-Intro-Deck.pptx — introduction deck with logo, summary, features.
Requires: pip install -r requirements-pptx.txt

Run from path-pulse-web:
  python generate-intro-pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BG = RGBColor(0x0B, 0x0E, 0x11)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
CYAN = RGBColor(0x00, 0xF5, 0xFF)
LIME = RGBColor(0x39, 0xFF, 0x14)


def _fill_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height).text_frame


def _para(tf, text, size_pt, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = align
    return p


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    root = Path(__file__).resolve().parent
    logo_path = root / "icon-512.png"

    # --- Slide 1: Title + logo (centered) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide)
    if logo_path.is_file():
        # Symmetrical: logo centered top
        slide.shapes.add_picture(str(logo_path), Inches(5.65), Inches(0.85), width=Inches(2.0), height=Inches(2.0))
    tf = _add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11.7), Inches(3.2))
    tf.word_wrap = True
    _para(tf, "TACTICAL BIO-LAB", 10, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    tf.paragraphs[0].space_after = Pt(4)
    _para(tf, "PATH-PULSE", 44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tf.paragraphs[-1].space_after = Pt(10)
    _para(
        tf,
        "Turn every walk into a mission. GPS routes, metabolic diagnostics, fuel tracking, "
        "and Ghost-Path racing—in your browser, Obsidian Lab aesthetic.",
        15,
        bold=False,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    tf.paragraphs[-1].line_spacing = 1.3

    # --- Slide 2: Two-column summary ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide)
    tf = _add_textbox(slide, Inches(0.65), Inches(0.45), Inches(12), Inches(0.9))
    _para(tf, "SUMMARY", 11, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    tf = _add_textbox(slide, Inches(0.65), Inches(0.95), Inches(12), Inches(0.7))
    _para(tf, "What is Path-Pulse?", 28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    w = Inches(5.9)
    left1 = Inches(0.65)
    left2 = Inches(6.78)
    top = Inches(2.0)
    h = Inches(4.2)

    tf1 = _add_textbox(slide, left1, top, w, h)
    tf1.word_wrap = True
    _para(tf1, "THE PRODUCT", 10, bold=True, color=LIME, align=PP_ALIGN.CENTER)
    tf1.paragraphs[0].space_after = Pt(8)
    _para(
        tf1,
        "A progressive web app that treats movement like tactical data: map expeditions, "
        "log fuel & steps, track BMI and weekly missions—without an app store install.",
        14,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    tf1.paragraphs[-1].line_spacing = 1.25

    tf2 = _add_textbox(slide, left2, top, w, h)
    tf2.word_wrap = True
    _para(tf2, "THE EXPERIENCE", 10, bold=True, color=LIME, align=PP_ALIGN.CENTER)
    tf2.paragraphs[0].space_after = Pt(8)
    _para(
        tf2,
        "Explorer's Oath → PRISM welcome → Home, Map, Report, Profile. Neon cyan & lime on "
        "obsidian—built for clarity and shareable lab vibes.",
        14,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    tf2.paragraphs[-1].line_spacing = 1.25

    # --- Slide 3: Features 1 (bullets) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide)
    tf = _add_textbox(slide, Inches(0.65), Inches(0.45), Inches(12), Inches(6.2))
    tf.word_wrap = True
    _para(tf, "FEATURES (1/2)", 11, bold=True, color=CYAN)
    _para(tf, "Core capabilities", 30, bold=True, color=WHITE)
    tf.paragraphs[-1].space_after = Pt(12)
    bullets = [
        "Map & expedition: Live GPS route, distance HUD, route replay & Ghost-Path.",
        "Home & missions: Kinetic step ring, weekly mission, explorer level & rank.",
        "PRISM: Onboard welcome after the Oath; optional TTS and haptics.",
        "Report & share: Weekly diagnostic, calories, share text or report image.",
    ]
    for b in bullets:
        p = _para(tf, "◆  " + b, 14, color=WHITE)
        p.space_after = Pt(8)
        p.line_spacing = 1.2

    # --- Slide 4: Features 2 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide)
    tf = _add_textbox(slide, Inches(0.65), Inches(0.45), Inches(12), Inches(6.2))
    tf.word_wrap = True
    _para(tf, "FEATURES (2/2)", 11, bold=True, color=CYAN)
    _para(tf, "Health & data", 30, bold=True, color=WHITE)
    tf.paragraphs[-1].space_after = Pt(12)
    bullets = [
        "Profile & lab: BMI, BMR, TDEE, body composition (Navy optional), progress calendar.",
        "Fuel gauge: Meals with timestamps, balance vs target, calorie→weight estimate.",
        "Your data: On-device storage, export JSON, PWA & offline map tiles.",
        "Accessibility: ARIA labels, keyboard nav, reduced-motion support.",
    ]
    for b in bullets:
        p = _para(tf, "◆  " + b, 14, color=WHITE)
        p.space_after = Pt(8)
        p.line_spacing = 1.2

    # --- Slide 5: Thank you ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide)
    if logo_path.is_file():
        slide.shapes.add_picture(str(logo_path), Inches(5.9), Inches(1.0), width=Inches(1.5), height=Inches(1.5))
    tf = _add_textbox(slide, Inches(0.8), Inches(2.85), Inches(11.7), Inches(3.5))
    tf.word_wrap = True
    _para(tf, "INITIALIZE THE GRID", 10, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    _para(tf, "Thank you", 36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tf.paragraphs[-1].space_after = Pt(10)
    _para(
        tf,
        "Open Path-Pulse in the browser, accept the Oath, and start your first expedition. "
        "Serve over HTTPS or localhost for GPS.",
        15,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )

    for i, s in enumerate(prs.slides):
        foot = s.shapes.add_textbox(
            Inches(0.65), prs.slide_height - Inches(0.5), Inches(12), Inches(0.35)
        ).text_frame
        fp = foot.paragraphs[0]
        fp.text = f"{i + 1} / {len(prs.slides)}"
        fp.font.size = Pt(10)
        fp.font.color.rgb = GRAY
        fp.alignment = PP_ALIGN.RIGHT

    out = root / "Path-Pulse-Intro-Deck.pptx"
    prs.save(out)
    print(f"Wrote: {out}")
    return out


if __name__ == "__main__":
    build()
