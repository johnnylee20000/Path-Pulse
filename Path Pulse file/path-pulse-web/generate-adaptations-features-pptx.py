"""
Path-Pulse — Adaptations & Features (full deck) + acknowledgements.

Generates: Path-Pulse-Adaptations-and-Features.pptx
- Symmetrical 16:9 layout (centered text, balanced margins)
- Introduction slide with animated Path-Pulse logo (GIF, ring spin like the web intro deck)
- Content slides + thank-you

Run from path-pulse-web:
  pip install -r requirements-pptx.txt
  python generate-adaptations-features-pptx.py
"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BG = RGBColor(0x0B, 0x0E, 0x11)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
CYAN = RGBColor(0x00, 0xF5, 0xFF)
LIME = RGBColor(0x39, 0xFF, 0x14)

# Match intro deck background (#0B0E11) and ring (#00F5FF)
_BG_RGB = (11, 14, 17)
_RING_RGB = (0, 245, 255)


def _fill_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _add_textbox(slide, left, top, width, height, valign=MSO_ANCHOR.TOP):
    shp = slide.shapes.add_textbox(left, top, width, height)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    return tf


def _para(tf, text, size_pt, bold=False, color=WHITE, align=PP_ALIGN.CENTER, space_after=Pt(6)):
    p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = align
    p.space_after = space_after
    return p


def build_animated_logo_gif(root: Path) -> Optional[Path]:
    """
    Build path-pulse-logo-animated.gif from icon-512.png — rotating cyan ring + centered logo
    (mirrors path-pulse-intro-deck.html ring-spin + logo).
    Returns path to GIF or None if skipped.
    """
    png = root / "icon-512.png"
    out = root / "path-pulse-logo-animated.gif"
    if not png.is_file():
        return None
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        return None

    size = 512
    frames = 36
    duration_ms = 55
    logo_raw = Image.open(png).convert("RGBA")
    # Logo sits inside the ring
    logo_max = int(size * 0.42)
    lw, lh = logo_raw.size
    scale = min(logo_max / lw, logo_max / lh)
    nw, nh = int(lw * scale), int(lh * scale)
    logo = logo_raw.resize((nw, nh), Image.Resampling.LANCZOS)

    cx, cy = size // 2, size // 2
    r_ring = int(size * 0.38)

    pil_frames = []
    for fi in range(frames):
        angle = int(fi * (360.0 / frames)) % 360
        im = Image.new("RGBA", (size, size), _BG_RGB + (255,))
        draw = ImageDraw.Draw(im)
        bbox = [cx - r_ring, cy - r_ring, cx + r_ring, cy + r_ring]
        # Rotating ring with a gap (reads like the web deck’s spinning ring)
        gap = 90
        a1 = angle
        a2 = angle + (360 - gap)
        # Pillow ImageDraw.arc uses fill= for stroke color (width>0)
        if a2 <= 360:
            draw.arc(bbox, a1, a2, fill=_RING_RGB, width=6)
        else:
            draw.arc(bbox, a1, 360, fill=_RING_RGB, width=6)
            draw.arc(bbox, 0, a2 - 360, fill=_RING_RGB, width=6)
        lx = (size - nw) // 2
        ly = (size - nh) // 2
        im.paste(logo, (lx, ly), logo)
        pil_frames.append(im)

    pil_frames[0].save(
        str(out),
        save_all=True,
        append_images=pil_frames[1:],
        duration=duration_ms,
        loop=0,
        optimize=True,
    )
    return out


def _add_intro_slide(prs, root: Path, total_slides: int):
    """Symmetrical title slide: animated GIF (or PNG) centered + introduction copy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide)
    sw, sh = prs.slide_width, prs.slide_height
    margin = Inches(0.85)
    content_w = sw - 2 * margin

    gif = build_animated_logo_gif(root)
    png = root / "icon-512.png"
    pic_w = Inches(2.35)
    left_pic = (sw - pic_w) / 2
    top_pic = Inches(0.55)
    try:
        if gif and gif.is_file():
            slide.shapes.add_picture(str(gif), left_pic, top_pic, width=pic_w, height=pic_w)
        elif png.is_file():
            slide.shapes.add_picture(str(png), left_pic, top_pic, width=pic_w, height=pic_w)
    except Exception:
        if png.is_file():
            try:
                slide.shapes.add_picture(str(png), left_pic, top_pic, width=pic_w, height=pic_w)
            except Exception:
                pass

    tf = _add_textbox(slide, margin, Inches(2.95), content_w, Inches(4.0))
    _para(tf, "INTRODUCTION", 11, bold=True, color=CYAN, align=PP_ALIGN.CENTER, space_after=Pt(8))
    _para(tf, "PATH-PULSE", 36, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=Pt(10))
    _para(
        tf,
        "Tactical Bio-Lab in the browser — movement as mission, data as protocol. "
        "This presentation walks through how Path-Pulse was adapted for the web, "
        "what each major feature does, and the reasoning behind it.",
        15,
        bold=False,
        color=GRAY,
        align=PP_ALIGN.CENTER,
        space_after=Pt(6),
    )
    tf.paragraphs[-1].line_spacing = 1.35

    foot = slide.shapes.add_textbox(margin, sh - Inches(0.5), content_w, Inches(0.35)).text_frame
    fp = foot.paragraphs[0]
    fp.text = f"1 / {total_slides}"
    fp.font.size = Pt(10)
    fp.font.color.rgb = GRAY
    fp.alignment = PP_ALIGN.CENTER


# Each entry: kicker, title, summary (one paragraph), logic (bullets — why / design rationale)
SLIDES = [
    {
        "kicker": "PATH-PULSE",
        "title": "Adaptations & Features",
        "summary": (
            "A tactical bio-lab in the browser: movement as mission, data as protocol. "
            "This deck summarizes how Path-Pulse was adapted for the web and why each major feature exists."
        ),
        "logic": [],
    },
    {
        "kicker": "ADAPTATION",
        "title": "Browser-first · PWA · Single codebase",
        "summary": (
            "Path-Pulse was adapted as a Progressive Web App so Explorers can run the Lab without an app store: "
            "HTTPS or localhost for GPS, install to home screen, data in localStorage by default."
        ),
        "logic": [
            "Logic: Remove install friction; keep privacy-first (on-device storage) until an optional API is added.",
            "One web codebase (path-pulse-web) keeps docs, marketing, and deployment aligned.",
        ],
    },
    {
        "kicker": "ONBOARDING",
        "title": "Explorer’s Oath & PRISM",
        "summary": (
            "The Oath frames consent and mindset; PRISM delivers a scripted welcome (voice optional) so the first "
            "moment feels like mission control—not a generic form."
        ),
        "logic": [
            "Logic: Narrative hooks retention; PRISM mirrors the product brief (Tactical Bio-Lab) and Flutter doc flow.",
            "Web Speech API + optional haptics approximate “assistant” without native AR.",
        ],
    },
    {
        "kicker": "PERSISTENCE",
        "title": "Service worker & offline map tiles",
        "summary": (
            "A service worker caches map tiles and core assets so the map degrades gracefully when connectivity drops."
        ),
        "logic": [
            "Logic: Expeditions happen outdoors—spotty signal should not blank the map entirely.",
        ],
    },
    {
        "kicker": "HOME",
        "title": "Kinetic ring · Mission · Widgets",
        "summary": (
            "Daily steps toward a goal, weekly distance mission (e.g. 2 km), explorer ID/level/rank—and optional Fuel "
            "and Expedition blocks you can hide for a cleaner HUD."
        ),
        "logic": [
            "Logic: Gamification without shame—objective progress bars; widgets optional so power users customize density.",
        ],
    },
    {
        "kicker": "MAP",
        "title": "Expedition · Replay · Ghost-Path",
        "summary": (
            "Leaflet + dark basemap: record a live GPS route, replay it as a time-lapse (5s/8s/15s, loop), and overlay "
            "your last completed route as a dashed “ghost” while a new expedition runs—race your past self."
        ),
        "logic": [
            "Logic: Ghost-Path brings “AR Ghost” narrative to the web map; last route persists in storage after reload.",
            "Replay turns a route into a shareable, cinematic recap.",
        ],
    },
    {
        "kicker": "TIME & STEPS",
        "title": "Steps & calories by calendar day",
        "summary": (
            "Steps and meals are keyed by local date; at midnight the “today” bucket resets so gauges and logs match "
            "how people think about their day—not a single endless counter."
        ),
        "logic": [
            "Logic: Aligns reporting with local sleep/wake cycles; avoids stale totals after crossing midnight.",
        ],
    },
    {
        "kicker": "FUEL",
        "title": "TDEE · Goals · Meals · Calorie→weight",
        "summary": (
            "BMR (Mifflin–St Jeor), TDEE = BMR × PAL from steps (WHO/FAO style). Choose Lose/Maintain/Gain to set the "
            "gauge target. Log breakfast/lunch/dinner with Save + timestamp. Balance vs burn; rough kg equivalent from kcal."
        ),
        "logic": [
            "Logic: International formulas; goal presets translate intent into a single daily target for the dial.",
            "Split meals improve recall vs one giant number; timestamps support habit auditing.",
        ],
    },
    {
        "kicker": "PROFILE",
        "title": "Units · BMI · Body composition · Calendar",
        "summary": (
            "Height (m/cm/ft), weight (kg/lbs), WHO BMI with optional Asian cut-offs. Body fat: Deurenberg or Navy "
            "(waist/neck/hip). Visceral proxy; muscle/bone estimates. Progress calendar shows steps, km, kcal, weight per day."
        ),
        "logic": [
            "Logic: Clinically familiar units; Navy improves fat estimate when circumferences exist.",
            "Calendar makes longitudinal progress visible without leaving Profile.",
        ],
    },
    {
        "kicker": "REPORT & DATA",
        "title": "Weekly diagnostic · Share · Export",
        "summary": (
            "Report: distance, steps, verdict, weekly calories. Share text includes calories + calorie goal. Export JSON "
            "for backup. Optional report image (html2canvas)."
        ),
        "logic": [
            "Logic: Shareable summaries for accountability; export supports data ownership and migration.",
        ],
    },
    {
        "kicker": "POLISH",
        "title": "Accessibility · Install · Portrait",
        "summary": (
            "ARIA labels, keyboard-friendly nav, reduced-motion respect. Dismissible install banner. Manifest portrait for "
            "phone-first use."
        ),
        "logic": [
            "Logic: Inclusive defaults; PWA hints reduce “just a website” friction on mobile.",
        ],
    },
    {
        "kicker": "ACKNOWLEDGEMENTS",
        "title": "With gratitude",
        "summary": (
            "Special thank you to the Ludwitt Academy, the Rose Foundation, the Hearts and Minds Program, and Cursor—"
            "for equipping me with knowledge, opportunity, and platforms to create; for helping expand my thinking; and for "
            "the tools that let me translate thoughts into actions and words. This work would not be the same without your support."
        ),
        "logic": [],
    },
]


def build():
    root = Path(__file__).resolve().parent
    total_slides = 1 + len(SLIDES)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    margin = Inches(0.85)
    content_w = prs.slide_width - 2 * margin

    _add_intro_slide(prs, root, total_slides)

    for i, data in enumerate(SLIDES):
        slide_idx = i + 2  # 1-based slide number after intro
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _fill_bg(slide)

        tf = _add_textbox(slide, margin, Inches(0.5), content_w, Inches(6.25))

        title_size = 28 if i == 0 else 26
        _para(tf, data["kicker"], 11, bold=True, color=CYAN, align=PP_ALIGN.CENTER, space_after=Pt(6))
        _para(tf, data["title"], title_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=Pt(12))

        _para(tf, data["summary"], 14, bold=False, color=GRAY, align=PP_ALIGN.CENTER, space_after=Pt(12))
        tf.paragraphs[-1].line_spacing = 1.28

        for line in data["logic"]:
            p = _para(tf, "◆  " + line, 13, bold=False, color=WHITE, align=PP_ALIGN.CENTER, space_after=Pt(8))
            p.line_spacing = 1.22

        foot = slide.shapes.add_textbox(
            margin, prs.slide_height - Inches(0.5), content_w, Inches(0.35)
        ).text_frame
        fp = foot.paragraphs[0]
        fp.text = f"{slide_idx} / {total_slides}"
        fp.font.size = Pt(10)
        fp.font.color.rgb = GRAY
        fp.alignment = PP_ALIGN.CENTER

    out = root / "Path-Pulse-Adaptations-and-Features.pptx"
    try:
        prs.save(out)
    except PermissionError:
        alt = root / "Path-Pulse-Adaptations-and-Features-generated.pptx"
        prs.save(alt)
        print(
            f"Could not overwrite {out.name} (close PowerPoint if it is open). "
            f"Saved as: {alt} ({total_slides} slides)"
        )
        out = alt
    else:
        print(f"Wrote: {out} ({total_slides} slides)")
    gif_path = root / "path-pulse-logo-animated.gif"
    if gif_path.is_file():
        print(f"Animated logo: {gif_path.name} (embedded on intro slide)")
    return out


if __name__ == "__main__":
    build()
