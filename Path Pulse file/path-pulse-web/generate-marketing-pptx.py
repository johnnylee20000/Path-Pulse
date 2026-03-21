"""
Build Path-Pulse-Marketing-Deck.pptx for Microsoft PowerPoint.
Requires: pip install python-pptx

Run from this folder:
  python generate-marketing-pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Path-Pulse theme
BG = RGBColor(0x0B, 0x0E, 0x11)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
CYAN = RGBColor(0x00, 0xF5, 0xFF)
SLIDES = [
    {
        "kicker": "TACTICAL BIO-LAB",
        "title": "Path-Pulse",
        "subtitle": (
            "Your body is the hardware. Movement is the mission. "
            "A browser-based fitness lab for steps, routes, and progress—no app store required."
        ),
        "bullets": [
            "Explorer-first onboarding with the Oath and PRISM welcome—then straight into the grid.",
            "Dark Obsidian UI with cyan & lime accents—built for focus, day or night.",
        ],
    },
    {
        "kicker": "ACCESS",
        "title": "Runs in the browser · PWA-ready",
        "subtitle": (
            "Open index.html via a local server for GPS, or install to your home screen "
            "for a standalone feel."
        ),
        "bullets": [
            "localStorage keeps your data on-device by default—private and offline-friendly.",
            "Optional Node backend for sync and Web Push when you set PATH_PULSE_API.",
        ],
    },
    {
        "kicker": "HOME",
        "title": "Daily pulse & goals",
        "subtitle": (
            "A kinetic step ring tracks progress toward your goal. See explorer ID, level, rank, "
            "BMI, BMR, burn, and protocol at a glance."
        ),
        "bullets": [
            "Manual step entry so you can align with your watch or phone totals.",
            "Quick actions—jump to Map, Report, Expedition, or scroll to Energy from the home strip.",
        ],
    },
    {
        "kicker": "MOTION",
        "title": "Auto steps from the device",
        "subtitle": (
            "Turn on motion to estimate steps from accelerometer / device motion—useful when "
            "you keep the app open while you move."
        ),
        "bullets": [
            "Smart messaging: expedition GPS distance doesn’t double-count against motion when both could apply.",
            "Tab visibility handling keeps step accounting consistent when you switch apps.",
        ],
    },
    {
        "kicker": "MAP",
        "title": "Expeditions on a live map",
        "subtitle": (
            "Leaflet map, geolocation, and Start / Stop Expedition to record your route in cyan. "
            "Distance HUD shows session track, odometer, and expedition step estimates."
        ),
        "bullets": [
            "Route replay with duration presets—time-lapse-style playback.",
            "Ghost-Path—race your previous route on the map.",
            "Demo routes for showcases when you’re testing without live GPS.",
        ],
    },
    {
        "kicker": "MISSION",
        "title": "Weekly mission & fuel",
        "subtitle": (
            "Walk a set distance each week (e.g. mission km), track completion, and monitor your "
            "energy / fuel widget alongside the rest of your stats."
        ),
        "bullets": [
            "Mission progress and status on Home—clear, gamified, and tied to your expedition distance.",
        ],
    },
    {
        "kicker": "REPORT",
        "title": "Diagnostics that respect your calendar",
        "subtitle": (
            "Day, week, month, and year views with charts and rollups—steps, distance, burn, "
            "water, exercise, and more."
        ),
        "bullets": [
            "Date keys use local calendar days so week and month boundaries match where you live—not UTC drift.",
            "Tap points for richer detail: meals, sleep, heart rate, blood pressure when logged.",
        ],
    },
    {
        "kicker": "PROFILE",
        "title": "Baseline & reminders",
        "subtitle": (
            "Weight, height, age, sex, units, body metrics, goals, and optional daily reminders—"
            "your lab profile drives BMR, BMI, and burn estimates."
        ),
        "bullets": [
            "Configurable calorie goal and predictive fueling mindset aligned with the app’s protocol language.",
        ],
    },
    {
        "kicker": "SHARE & SYNC",
        "title": "Capture & optional cloud",
        "subtitle": (
            "html2canvas support for sharing visuals. Connect the optional API for device sync "
            "and Web Push nudges when the backend is configured."
        ),
        "bullets": [
            "See backend/README.md for VAPID keys and deployment notes.",
        ],
    },
    {
        "kicker": "START",
        "title": "Initialize the grid",
        "subtitle": (
            "Path-Pulse is your tactical bio-lab for movement data: oath, map, report, profile—"
            "one tab, full mission."
        ),
        "bullets": [
            "Open index.html in the browser (after starting a local server).",
            "Tip: serve over localhost or HTTPS for best geolocation results.",
        ],
    },
]


def _fill_slide_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height).text_frame


def _para(tf, text, size_pt, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = align
    return p


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    margin = Inches(0.65)
    content_w = prs.slide_width - 2 * margin

    for i, data in enumerate(SLIDES):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        _fill_slide_background(slide)

        tf = _add_textbox(slide, margin, Inches(0.45), content_w, Inches(6.5))
        tf.word_wrap = True

        _para(tf, data["kicker"], 11, bold=True, color=CYAN)
        tf.paragraphs[0].space_after = Pt(6)

        title_size = 40 if i == 0 else 32
        _para(tf, data["title"], title_size, bold=True, color=WHITE)
        tf.paragraphs[-1].space_after = Pt(10)

        _para(tf, data["subtitle"], 15, bold=False, color=GRAY)
        tf.paragraphs[-1].space_after = Pt(14)
        tf.paragraphs[-1].line_spacing = 1.25

        for b in data["bullets"]:
            p = _para(tf, "◆  " + b, 14, bold=False, color=WHITE)
            p.level = 0
            p.space_after = Pt(8)
            p.line_spacing = 1.2
            # Diamond color hint: use default white; optional bullet styling needs XML tweak

        # Footer: slide number
        foot = slide.shapes.add_textbox(
            margin, prs.slide_height - Inches(0.55), content_w, Inches(0.35)
        ).text_frame
        fp = foot.paragraphs[0]
        fp.text = f"{i + 1} / {len(SLIDES)}"
        fp.font.size = Pt(10)
        fp.font.color.rgb = GRAY
        fp.alignment = PP_ALIGN.RIGHT

    out = Path(__file__).resolve().parent / "Path-Pulse-Marketing-Deck.pptx"
    prs.save(out)
    print(f"Wrote: {out}")
    return out


if __name__ == "__main__":
    build()
